import aiohttp
import asyncio
import os
import re
from bs4 import BeautifulSoup
from urllib.parse import urljoin
from typing import Dict, List, Optional
import logging
from pptx import Presentation

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def extract_text_from_ppt(ppt_path):
    """从PPT提取文本"""
    if not os.path.exists(ppt_path):
        return f"错误：文件 '{ppt_path}' 不存在"
    try:
        prs = Presentation(ppt_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text.strip())
        return "\n".join(all_text)
    except Exception as e:
        return f"读取PPT文件时出错: {str(e)}"


async def extract_content_from_url(url: str) -> Dict:
    """异步提取网页内容（文本+图片URL），当内容为空时自动重试两次"""
    max_attempts = 3  # 最多尝试3次（1次初始+2次重试）
    attempt = 0

    while attempt < max_attempts:
        attempt += 1
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
                'Accept-Language': 'zh-CN,zh;q=0.8,zh-TW;q=0.7,zh-HK;q=0.5,en-US;q=0.3,en;q=0.2',
            }

            logger.info(f"第 {attempt}/{max_attempts} 次尝试访问页面: {url}")
            async with aiohttp.ClientSession() as session:
                async with session.get(url, headers=headers, timeout=15) as response:
                    if response.status != 200:
                        result = {"document": f"请求失败，状态码: {response.status}", "image_urls": []}
                    else:
                        html_content = await response.text()
                        # 提取文本和图片URL（按网页中出现的顺序）
                        document, image_urls = await extract_text_and_images_in_order(html_content, url)
                        result = {
                            "document": document,
                            "image_urls": image_urls
                        }

            # 检查是否需要重试（document为空或仅含错误信息）
            if not result["document"].strip() or \
                    result["document"].startswith(("请求失败", "获取网页内容时出错", "提取内容时出错")):
                if attempt < max_attempts:
                    logger.warning(f"第 {attempt} 次爬取失败（内容为空或错误），准备重试...")
                    await asyncio.sleep(1)  # 等待1秒后重试，避免频繁请求
                    continue  # 继续下一次尝试
                else:
                    logger.error(f"已达到最大重试次数（{max_attempts}次），爬取仍失败")

            # 若无需重试，直接返回结果
            return result

        except Exception as e:
            error_msg = f"第 {attempt} 次爬取发生异常: {str(e)}"
            logger.error(error_msg)
            if attempt < max_attempts:
                logger.warning("准备重试...")
                await asyncio.sleep(1)
            else:
                return {"document": error_msg, "image_urls": []}

    # 所有尝试失败后返回最终错误
    return {"document": f"超过最大重试次数（{max_attempts}次），无法获取网页内容", "image_urls": []}


async def extract_text_and_images_in_order(html_content: str, base_url: str) -> (str, List[str]):
    """按网页中出现的顺序提取文本和图片URL，整合为document"""
    try:
        # 同步解析HTML（使用线程池避免阻塞事件循环）
        soup = await asyncio.to_thread(BeautifulSoup, html_content, 'html.parser')

        # 移除无关标签
        for tag in soup(["script", "style", "noscript", "meta", "link"]):
            tag.decompose()

        # 按顺序遍历所有有意义的节点（文本节点和图片节点）
        document_parts = []  # 存储按顺序排列的文本和图片URL
        image_urls = []  # 单独存储图片URL（用于后续引用）

        # 遍历<body>下的所有直接子节点，按顺序提取内容
        for child in soup.body.children if soup.body else []:
            if child.name is None:
                # 文本节点：清理后添加到document
                text = child.strip()
                if text:
                    cleaned_text = remove_duplicates_and_noise(text)
                    if cleaned_text:
                        document_parts.append(f"[文本内容]\n{cleaned_text}\n")
            elif child.name == 'img':
                # 图片节点：提取URL并添加到document
                src = child.get('src') or child.get('data-src')
                if src and not src.startswith('data:image'):
                    # 补全相对URL
                    img_url = urljoin(base_url, src) if not src.startswith(('http://', 'https://')) else src
                    document_parts.append(f"[图片内容]\nURL: {img_url}\n")
                    image_urls.append(img_url)
            else:
                # 其他标签：递归提取内部文本（如p、div等）
                inner_text = child.get_text(strip=True, separator='\n')
                if inner_text:
                    cleaned_text = remove_duplicates_and_noise(inner_text)
                    if cleaned_text:
                        document_parts.append(f"[文本内容]\n{cleaned_text}\n")

        # 合并所有部分为完整document（限制长度避免溢出）
        document = ''.join(document_parts)
        if len(document) > 10000:
            document = document[:10000] + "\n[内容过长，已截断]"

        return document, image_urls

    except Exception as e:
        logger.error(f"提取文本和图片顺序时出错: {str(e)}")
        return f"提取内容时出错: {str(e)}", []


def remove_duplicates_and_noise(text):
    """清理文本中的重复内容和噪声"""
    lines = text.split('\n')
    lines = [line for line in lines if line.strip() and len(line.strip()) > 5]
    seen = set()
    unique_lines = [line for line in lines if line not in seen and not seen.add(line)]

    # 过滤广告、平台信息等噪声
    noise_patterns = [r'小红书.*号.*', r'©.*小红书', r'下载小红书', r'APP.*内打开', r'广告']
    filtered_lines = []
    for line in unique_lines:
        if not any(re.search(pattern, line) for pattern in noise_patterns):
            filtered_lines.append(line)
    return '\n'.join(filtered_lines)


def read_text_file(file_path):
    """读取文本文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"读取文件时出错: {str(e)}"
